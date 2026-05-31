"""
==============================================================================
Service: Document Generator — v2.3 P1 文書簡報 (.pptx) 渲染
         Document Generator — v2.3 P1 presentation (.pptx) rendering
==============================================================================
ZH: 用途：把 AI 產出的結構化 spec（dict）渲染成 .pptx，並送進該使用者的
    code-server 容器 /home/coder/outputs/。

    為什麼送進容器：job-scheduler 與學生容器是不同 container，學生在
    code-server (VS Code in Browser) 看到的檔案系統是 home_<user_id> volume
    掛在 /home/coder。job-scheduler 透過 Docker SDK（lab_manager 已有 client）
    用 put_archive 把生成的 .pptx 拷貝進 running 容器，學生即可在 outputs/
    直接看到並下載。

EN: Purpose: Render the AI-produced structured spec (dict) into a .pptx and
    deliver it into the user's code-server container at /home/coder/outputs/.

    Why into the container: job-scheduler and the student container are
    separate. The student's code-server sees home_<user_id> volume mounted at
    /home/coder. job-scheduler copies the generated .pptx into the running
    container via the Docker SDK (put_archive), so the student finds it under
    outputs/ immediately.
==============================================================================
"""

from __future__ import annotations

import io
import logging
import tarfile
import time
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Pt
from docker.errors import NotFound, APIError

from . import lab_manager

logger = logging.getLogger(__name__)

OUTPUTS_DIR = "/home/coder/outputs"
_MAX_SLIDES = 40          # ZH: PoC 上限，避免 AI 失控生 200 張 | EN: PoC cap
_MAX_BULLETS = 12         # ZH: 每張 bullet 上限 | EN: per-slide bullet cap


class DocumentGenerationError(Exception):
    """ZH: 生成流程可預期的失敗（回給使用者友善訊息）| EN: Expected failure"""


# ==============================================================================
# ZH: spec 正規化與驗證 | EN: spec normalization & validation
# ==============================================================================

def _clean_text(value, fallback: str = "") -> str:
    if value is None:
        return fallback
    text = str(value).strip()
    return text or fallback


def _normalize_spec(spec: dict) -> dict:
    """ZH: 容錯地把 AI spec 整理成可渲染結構 | EN: Tolerant normalization."""
    if not isinstance(spec, dict):
        raise DocumentGenerationError("spec 不是物件 (expected JSON object)")

    title = _clean_text(spec.get("title"), "未命名簡報")
    subtitle = _clean_text(spec.get("subtitle"))

    raw_slides = spec.get("slides")
    if not isinstance(raw_slides, list) or not raw_slides:
        raise DocumentGenerationError("spec.slides 缺少或為空 (need at least 1 slide)")

    slides: list[dict] = []
    for item in raw_slides[:_MAX_SLIDES]:
        if not isinstance(item, dict):
            continue
        s_title = _clean_text(item.get("title"), "（無標題）")
        raw_bullets = item.get("bullets")
        bullets: list[str] = []
        if isinstance(raw_bullets, list):
            for b in raw_bullets[:_MAX_BULLETS]:
                bt = _clean_text(b)
                if bt:
                    bullets.append(bt)
        notes = _clean_text(item.get("notes"))
        slides.append({"title": s_title, "bullets": bullets, "notes": notes})

    if not slides:
        raise DocumentGenerationError("spec.slides 內無有效投影片 (no valid slide)")

    return {"title": title, "subtitle": subtitle, "slides": slides}


# ==============================================================================
# ZH: python-pptx 渲染 | EN: python-pptx rendering
# ==============================================================================

def _render_pptx_bytes(norm: dict) -> bytes:
    """ZH: 用內建 template 渲染成 .pptx，回傳 bytes | EN: Render to .pptx bytes."""
    prs = Presentation()  # ZH: 預設 16:9-ish template | EN: built-in default template

    # --- 封面 / Title slide ---
    title_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = norm["title"]
    if cover.placeholders and len(cover.placeholders) > 1:
        try:
            cover.placeholders[1].text = norm["subtitle"] or ""
        except (KeyError, IndexError):
            pass

    # --- 內容頁 / Content slides (layout 1 = Title and Content) ---
    content_layout = prs.slide_layouts[1]
    for s in norm["slides"]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = s["title"]

        # ZH: 取內容 placeholder（idx 1）寫 bullet | EN: body placeholder
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is not None:
            tf = body.text_frame
            tf.clear()
            if s["bullets"]:
                for i, bullet in enumerate(s["bullets"]):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = bullet
                    para.level = 0
                    para.font.size = Pt(18)
            else:
                tf.paragraphs[0].text = ""

        # ZH: 備忘稿 | EN: speaker notes
        if s["notes"]:
            slide.notes_slide.notes_text_frame.text = s["notes"]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ==============================================================================
# ZH: 送進學生容器 | EN: deliver into the student's container
# ==============================================================================

def _safe_filename(title: str) -> str:
    base = "".join(c for c in title if c.isalnum() or c in (" ", "-", "_")).strip()
    base = base.replace(" ", "_")[:40] or "presentation"
    ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    return f"{base}_{ts}.pptx"


def _put_into_container(user_id: str, filename: str, data: bytes) -> str:
    """
    ZH: 把 bytes 以 put_archive 拷貝進 running 容器的 OUTPUTS_DIR。
        回傳容器內完整路徑。容器不存在/未啟動 → DocumentGenerationError。
    EN: Copy bytes into running container's OUTPUTS_DIR via put_archive.
    """
    lc = lab_manager.get_lifecycle()
    container_name = lc._container_name(user_id)
    try:
        container = lc.client.containers.get(container_name)
    except NotFound:
        raise DocumentGenerationError(
            "你的 Notebook 容器尚未啟動，請先到「Notebook」分頁啟動，再重新生成。"
            " | Your Notebook container is not running; start it first."
        )

    if container.status != "running":
        raise DocumentGenerationError(
            "你的 Notebook 容器目前未在執行，請先啟動再生成。"
            " | Your Notebook container is not running."
        )

    # ZH: 確保 outputs 目錄存在（volume 掛載可能蓋掉 image 內的目錄）
    # EN: Ensure outputs dir exists (volume mount may shadow image dir)
    try:
        container.exec_run(f"mkdir -p {OUTPUTS_DIR}", user="coder")
    except APIError as e:
        logger.warning("mkdir outputs failed (continuing): %s", e)

    # ZH: 打包單一檔案成 tar，put_archive 解包到 OUTPUTS_DIR
    # EN: Pack single file into a tar, put_archive extracts into OUTPUTS_DIR
    tar_stream = io.BytesIO()
    with tarfile.open(fileobj=tar_stream, mode="w") as tar:
        info = tarfile.TarInfo(name=filename)
        info.size = len(data)
        info.mtime = int(time.time())
        info.mode = 0o644
        tar.addfile(info, io.BytesIO(data))
    tar_stream.seek(0)

    ok = container.put_archive(OUTPUTS_DIR, tar_stream.getvalue())
    if not ok:
        raise DocumentGenerationError("檔案寫入容器失敗 | Failed to write file into container")

    # ZH: 修正擁有者為 coder（put_archive 以 root 解包）
    # EN: Fix ownership to coder (put_archive extracts as root)
    try:
        container.exec_run(f"chown coder:coder {OUTPUTS_DIR}/{filename}", user="root")
    except APIError:
        pass

    return f"{OUTPUTS_DIR}/{filename}"


# ==============================================================================
# ZH: 對外主函式 | EN: public entry
# ==============================================================================

def generate_presentation(spec: dict, user_id: str) -> dict:
    """
    ZH: 由 chat.py presentation 分支呼叫。
        成功 → {"ok": True, "filename", "path", "slides"}
        失敗 → {"ok": False, "error": "<友善訊息>"}
    EN: Called by chat.py presentation branch.
    """
    try:
        norm = _normalize_spec(spec)
        data = _render_pptx_bytes(norm)
        filename = _safe_filename(norm["title"])
        path = _put_into_container(user_id, filename, data)
        logger.info(
            "Generated presentation for user %s: %s (%d slides, %d bytes)",
            user_id[:8], filename, len(norm["slides"]), len(data),
        )
        return {
            "ok": True,
            "filename": filename,
            "path": path,
            "slides": len(norm["slides"]),
        }
    except DocumentGenerationError as e:
        logger.warning("Presentation generation failed for %s: %s", user_id[:8], e)
        return {"ok": False, "error": str(e)}
    except Exception as e:  # noqa: BLE001 — PoC：任何渲染例外都回友善訊息
        logger.error("Unexpected presentation error for %s: %s", user_id[:8], e, exc_info=True)
        return {"ok": False, "error": f"簡報生成發生未預期錯誤 | Unexpected error: {e}"}
